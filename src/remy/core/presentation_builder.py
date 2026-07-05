"""
PPTX Presentation Generator for Remy Agent.

Generates professional PowerPoint presentations from structured data.
Uses python-pptx with full Unicode/Cyrillic support.
"""

import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)


# ============================================================
# COLOR SCHEME (matches report_builder.py)
# ============================================================

COLORS = {
    "primary": RGBColor(0x1A, 0x1A, 0x2E),
    "secondary": RGBColor(0x16, 0x21, 0x3E),
    "accent": RGBColor(0x0F, 0x34, 0x60),
    "highlight": RGBColor(0xE9, 0x45, 0x60),
    "text": RGBColor(0x2D, 0x2D, 0x2D),
    "text_light": RGBColor(0x99, 0x99, 0x99),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg_light": RGBColor(0xF5, 0xF5, 0xF5),
    "trust_high": RGBColor(0x27, 0xAE, 0x60),
    "trust_mid": RGBColor(0xF3, 0x9C, 0x12),
    "trust_low": RGBColor(0xE7, 0x4C, 0x3C),
}


# ============================================================
# PRESENTATION BUILDER
# ============================================================

class PresentationBuilder:
    """
    Builds a PPTX presentation step-by-step.

    Usage:
        pres = PresentationBuilder("Research Report", "AI Agent Analysis")
        pres.add_section("Introduction", "This report covers...")
        pres.add_bullets("Key Points", ["Point 1", "Point 2"])
        pres.add_table("Data", headers, rows)
        pres.save("report.pptx")
    """

    def __init__(
        self,
        title: str,
        subtitle: str = "",
        author: str = "Remy AI Agent",
        output_dir: Optional[str] = None,
    ):
        self.title = title
        self.subtitle = subtitle
        self.author = author
        self.output_dir = Path(output_dir or "./data/presentations")
        self.output_dir.mkdir(parents=True, exist_ok=True)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)
        self._build_cover()

    # ------ helpers ------

    def _add_bg(self, slide, color=None):
        """Fill slide background with a solid color."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color or COLORS["white"]

    def _add_footer(self, slide, text: str = ""):
        """Add a subtle footer bar at the bottom of a slide."""
        footer_text = text or f"{self.author}  |  {datetime.now().strftime('%Y-%m-%d')}"
        left = Inches(0)
        top = self.prs.slide_height - Inches(0.45)
        width = self.prs.slide_width
        height = Inches(0.45)

        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["primary"]
        shape.line.fill.background()

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = footer_text
        p.font.size = Pt(10)
        p.font.color.rgb = COLORS["text_light"]
        p.alignment = PP_ALIGN.CENTER

    def _set_text(self, text_frame, text: str, font_size: int = 18,
                  color=None, bold: bool = False, alignment=None):
        """Set text in a text frame with styling."""
        text_frame.word_wrap = True
        p = text_frame.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color or COLORS["text"]
        p.font.bold = bold
        if alignment:
            p.alignment = alignment

    def _add_accent_bar(self, slide, left, top, width=Inches(0.08), height=Inches(0.6)):
        """Add a vertical accent bar."""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["highlight"]
        shape.line.fill.background()
        return shape

    # ------ slide types ------

    def _build_cover(self):
        """Title slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # blank
        self._add_bg(slide, COLORS["primary"])

        # Accent stripe
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.8),
            self.prs.slide_width, Inches(0.06),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["highlight"]
        shape.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
        self._set_text(txBox.text_frame, self.title, font_size=40,
                       color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)

        # Subtitle
        if self.subtitle:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11), Inches(1))
            self._set_text(txBox.text_frame, self.subtitle, font_size=22,
                           color=COLORS["text_light"], alignment=PP_ALIGN.CENTER)

        # Author + date
        now = datetime.now().strftime("%d %B %Y")
        txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.6))
        self._set_text(txBox.text_frame, f"{self.author}  •  {now}",
                       font_size=14, color=COLORS["text_light"], alignment=PP_ALIGN.CENTER)

    def add_section(self, title: str, body: str = ""):
        """Add a content slide with title and body text."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_footer(slide)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
        self._set_text(txBox.text_frame, title, font_size=28,
                       color=COLORS["primary"], bold=True)

        # Accent underline
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.25),
            Inches(2), Inches(0.04),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["highlight"]
        shape.line.fill.background()

        # Body
        if body:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            # Split on double-newline for paragraphs
            paragraphs = body.split("\n\n") if "\n\n" in body else body.split("\n")
            for i, para_text in enumerate(paragraphs):
                para_text = para_text.strip()
                if not para_text:
                    continue
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para_text
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS["text"]
                p.space_after = Pt(10)

    def add_subsection(self, title: str, body: str = ""):
        """Add a subsection slide (lighter title style)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_footer(slide)

        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
        self._set_text(txBox.text_frame, title, font_size=24,
                       color=COLORS["secondary"], bold=True)

        if body:
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.5), Inches(5.2))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, para_text in enumerate(body.split("\n\n") if "\n\n" in body else body.split("\n")):
                para_text = para_text.strip()
                if not para_text:
                    continue
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = para_text
                p.font.size = Pt(16)
                p.font.color.rgb = COLORS["text"]
                p.space_after = Pt(8)

    def add_bullets(self, title: str, items: list[str]):
        """Add a slide with title and bullet points."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_footer(slide)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
        self._set_text(txBox.text_frame, title, font_size=28,
                       color=COLORS["primary"], bold=True)

        # Accent underline
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.25),
            Inches(2), Inches(0.04),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["highlight"]
        shape.line.fill.background()

        # Bullets with accent bars
        y_pos = Inches(1.6)
        for item in items:
            self._add_accent_bar(slide, Inches(0.8), y_pos, height=Inches(0.45))

            txBox = slide.shapes.add_textbox(Inches(1.1), y_pos, Inches(11), Inches(0.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(16)
            p.font.color.rgb = COLORS["text"]

            y_pos += Inches(0.6)
            if y_pos > Inches(6.5):
                break  # prevent overflow

    def add_quote(self, text: str, author: str = ""):
        """Add a quote slide."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, COLORS["primary"])

        # Large accent bar
        self._add_accent_bar(slide, Inches(1.5), Inches(2), height=Inches(2.5))

        # Quote text
        txBox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(9), Inches(2.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f'"{text}"'
        p.font.size = Pt(22)
        p.font.color.rgb = COLORS["white"]
        p.font.italic = True

        if author:
            p = tf.add_paragraph()
            p.text = f"— {author}"
            p.font.size = Pt(14)
            p.font.color.rgb = COLORS["text_light"]
            p.space_before = Pt(20)

    def add_table(self, title: str, headers: list[str], rows: list[list[str]]):
        """Add a slide with a styled data table."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide)
        self._add_footer(slide)

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.5), Inches(0.9))
        self._set_text(txBox.text_frame, title, font_size=24,
                       color=COLORS["primary"], bold=True)

        n_cols = len(headers)
        n_rows = len(rows) + 1  # +1 for header
        # Limit rows to prevent overflow
        max_rows = min(n_rows, 12)

        table_width = Inches(11.5)
        table_height = Inches(0.4) * max_rows
        left = Inches(0.8)
        top = Inches(1.5)

        table_shape = slide.shapes.add_table(max_rows, n_cols, left, top, table_width, table_height)
        table = table_shape.table

        # Column widths
        col_w = int(table_width / n_cols)
        for i in range(n_cols):
            table.columns[i].width = col_w

        # Header row
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["primary"]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = COLORS["white"]
            p.font.bold = True
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Data rows
        for r_idx, row in enumerate(rows[:max_rows - 1]):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx + 1, c_idx)
                cell.text = str(val)
                if r_idx % 2 == 1:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLORS["bg_light"]
                p = cell.text_frame.paragraphs[0]
                p.font.size = Pt(11)
                p.font.color.rgb = COLORS["text"]
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    def add_section_divider(self, title: str):
        """Add a section divider slide (dark background, large text)."""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, COLORS["secondary"])

        # Accent stripe
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.5),
            self.prs.slide_width, Inches(0.05),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["highlight"]
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.2))
        self._set_text(txBox.text_frame, title, font_size=36,
                       color=COLORS["white"], bold=True, alignment=PP_ALIGN.CENTER)

    def save(self, filename: Optional[str] = None) -> str:
        """Save PPTX and return the filepath."""
        if not filename:
            safe_title = "".join(
                c if c.isalnum() or c in " -_" else "_"
                for c in self.title
            ).strip().replace(" ", "_")
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{safe_title}_{timestamp}.pptx"

        filepath = self.output_dir / filename
        self.prs.save(str(filepath))
        return str(filepath)
